#!/usr/bin/env python3
import json
import os
from pptx import Presentation

BASE_DIR = 'Z:/projects/goose'
OUTPUT_DIR = BASE_DIR + '/AI training/output'
CONTENT_DIR = BASE_DIR + '/AI training/ppt-content'

FILES = [
    {'file': OUTPUT_DIR + '/AI虛擬組織_Executive版.pptx', 'label': 'Executive版', 'expected_pages': 15, 'mid_page': 5},
    {'file': OUTPUT_DIR + '/AI虛擬組織_Developer版.pptx', 'label': 'Developer版', 'expected_pages': 30, 'mid_page': 5},
    {'file': OUTPUT_DIR + '/AI虛擬組織_PowerUser版.pptx', 'label': 'PowerUser版', 'expected_pages': 20, 'mid_page': 5},
]

EMU = 914400
EXP_W = int(10 * EMU)
EXP_H = int(7.5 * EMU)

def get_text(slide):
    t = []
    for s in slide.shapes:
        if s.has_text_frame:
            for p in s.text_frame.paragraphs:
                line = p.text.strip()
                if line: t.append(line)
    return t

def get_notes(slide):
    ns = slide.notes_slide
    if ns is None: return ''
    ntf = ns.notes_text_frame
    return ntf.text.strip() if ntf else ''

def has_calibri(prs):
    for i in range(min(3, len(prs.slides))):
        for sh in prs.slides[i].shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name and 'Calibri' in r.font.name:
                            return True
    return False

def check(entry):
    r = {'label': entry['label'], 'bugs': []}
    fp = entry['file']
    if not os.path.exists(fp):
        r['bugs'].append({'page': 'N/A', 'issue': 'File not found', 'severity': 'critical'})
        return r
    try:
        prs = Presentation(fp)
    except Exception as e:
        r['bugs'].append({'page': 'N/A', 'issue': str(e), 'severity': 'critical'})
        return r
    n = len(prs.slides)
    r['page_count'] = n
    r['page_count_ok'] = (n == entry['expected_pages'])
    if not r['page_count_ok']:
        r['bugs'].append({'page': 'N/A', 'issue': f"頁數不符: 預期 {entry['expected_pages']}, 實際 {n}", 'severity': 'critical'})
    w, h = prs.slide_width, prs.slide_height
    r['size_actual'] = f'{w/EMU:.2f}x{h/EMU:.2f}'
    r['size_ok'] = (abs(w - EXP_W) < 10000 and abs(h - EXP_H) < 10000)
    if not r['size_ok']:
        r['bugs'].append({'page': 'N/A', 'issue': f'尺寸不符: {r["size_actual"]}', 'severity': 'major'})
    ct = get_text(prs.slides[0])
    r['cover_texts'] = ct[:3]
    r['cover_content_ok'] = bool(ct)
    if not r['cover_content_ok']:
        r['bugs'].append({'page': 1, 'issue': '封面無文字', 'severity': 'critical'})
    lt = get_text(prs.slides[-1])
    r['last_page_texts'] = lt[:3]
    r['last_page_content_ok'] = bool(lt)
    if not r['last_page_content_ok']:
        r['bugs'].append({'page': n, 'issue': '末頁無文字', 'severity': 'major'})
    nc = sum(1 for i in range(min(5, n)) if get_notes(prs.slides[i]))
    r['notes_count'] = nc
    r['notes_ok'] = (nc >= 2)
    if not r['notes_ok']:
        r['bugs'].append({'page': '1-5', 'issue': f'Notes 不足: {nc} 頁', 'severity': 'major'})
    mid_i = entry['mid_page'] - 1
    if mid_i < n:
        mt = get_text(prs.slides[mid_i])
        r['mid_page_texts'] = mt[:3]
        r['mid_page_ok'] = (len(mt) >= 2)
        if not r['mid_page_ok']:
            r['bugs'].append({'page': entry['mid_page'], 'issue': f'第{entry["mid_page"]}頁內容過少', 'severity': 'minor'})
    r['calibri_found'] = has_calibri(prs)
    if not r['calibri_found']:
        r['bugs'].append({'page': '前3頁', 'issue': '未偵測 Calibri 字體', 'severity': 'minor'})
    return r

def verdict(r):
    sev = {b['severity'] for b in r['bugs']}
    if 'critical' in sev or 'major' in sev: return 'FAIL'
    return 'PASS'

results = [check(e) for e in FILES]
for r in results:
    v = verdict(r)
    print(f"
{'='*50}")
    print(f"檢查: {r['label']}")
    print(f"  頁數: {r.get('page_count','?')} {'OK' if r.get('page_count_ok') else 'FAIL'}")
    print(f"  尺寸: {r.get('size_actual','?')} {'OK' if r.get('size_ok') else 'FAIL'}")
    print(f"  封面: {'OK' if r.get('cover_content_ok') else 'FAIL'} {r.get('cover_texts',[][:2])}")
    print(f"  未頁: {'OK' if r.get('last_page_content_ok') else 'FAIL'} {r.get('last_page_texts',[][:2])}")
    print(f"  Notes: {r.get('notes_count','?')} 頁 {'OK' if r.get('notes_ok') else 'FAIL'}")
    print(f"  第5頁: {'OK' if r.get('mid_page_ok') else 'FAIL'} {r.get('mid_page_texts',[][:2])}")
    print(f"  Calibri: {'found' if r.get('calibri_found') else 'not detected'}")
    for b in r['bugs']:
        print(f"  [{b['severity'].upper()}] p{b['page']}: {b['issue']}")
    print(f"  裁定: {v}")

with open(BASE_DIR + '/.claude/scripts/qa-check-ppt-result.json', 'w', encoding='utf-8') as f:
    json.dump([{'label': r['label'], 'page_count': r.get('page_count'), 'size_actual': r.get('size_actual'), 'size_ok': r.get('size_ok'), 'cover_ok': r.get('cover_content_ok'), 'notes_count': r.get('notes_count'), 'notes_ok': r.get('notes_ok'), 'mid_page_ok': r.get('mid_page_ok'), 'calibri_found': r.get('calibri_found'), 'bugs': r['bugs'], 'verdict': verdict(r)} for r in results], f, ensure_ascii=False, indent=2)
print('
JSON written.')
