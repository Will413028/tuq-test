#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generic JSON slides → PPTX converter.
Usage: python gen-pptx-from-json.py <input.json> <output.pptx>
"""

import json
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Brand colors
DARK_BLUE = RGBColor(30, 39, 97)       # #1E2761
ICE_BLUE = RGBColor(202, 220, 252)     # #CADCFC
WHITE = RGBColor(255, 255, 255)        # #FFFFFF
DARK_GRAY = RGBColor(44, 62, 80)       # #2C3E50
ACCENT_BLUE = RGBColor(0, 102, 204)    # #0066CC
LIGHT_BG = RGBColor(240, 244, 255)     # #F0F4FF
CODE_BG = RGBColor(26, 26, 46)         # #1a1a2e

def hex_to_rgb(hex_str):
    """Convert hex color string to RGBColor."""
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                font_color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri', word_wrap=True):
    """Add a textbox with text to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_points(slide, left, top, width, height, items, font_size=14,
                      font_color=DARK_GRAY, font_name='Calibri', spacing=1.2):
    """Add bullet point list to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        text = str(item)
        # Handle nested items
        if isinstance(item, dict):
            text = item.get('text', item.get('title', str(item)))

        p.text = f"• {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = Pt(font_size * spacing * 0.5)

    return txBox

def add_title_bar(slide, title, is_dark_bg=False):
    """Add a branded title bar at the top of a slide."""
    if is_dark_bg:
        # White title on dark background
        add_textbox(slide, 0.5, 0.4, 9, 0.8, title,
                   font_size=36, font_color=WHITE, bold=True,
                   alignment=PP_ALIGN.LEFT)
        # Accent line
        line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.2), Inches(2), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = ICE_BLUE
        line.line.fill.background()
    else:
        # Dark blue title bar
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(10), Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_BLUE
        bar.line.fill.background()

        add_textbox(slide, 0.5, 0.15, 9, 0.8, title,
                   font_size=32, font_color=WHITE, bold=True)

def add_notes(slide, notes_text):
    """Add speaker notes to a slide."""
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = str(notes_text)

def flatten_content(content):
    """Extract text items from various content structures."""
    if isinstance(content, str):
        return [content]
    if isinstance(content, list):
        result = []
        for item in content:
            if isinstance(item, str):
                result.append(item)
            elif isinstance(item, dict):
                # Try common keys
                for key in ['text', 'title', 'point', 'description', 'label', 'name']:
                    if key in item:
                        val = item[key]
                        # Add sub-items if present
                        desc = item.get('description', item.get('detail', item.get('details', '')))
                        if desc and key != 'description':
                            result.append(f"{val} — {desc}")
                        else:
                            result.append(str(val))
                        break
                else:
                    result.append(str(item))
            else:
                result.append(str(item))
        return result
    if isinstance(content, dict):
        # Try to extract meaningful content from dict
        items = []
        for key, val in content.items():
            if key in ('style', 'layout_type', 'slide_number', 'title', 'subtitle', 'notes'):
                continue
            if isinstance(val, list):
                items.extend(flatten_content(val))
            elif isinstance(val, str):
                items.append(f"{key}: {val}")
            elif isinstance(val, dict):
                items.extend(flatten_content(val))
        return items
    return [str(content)]

def build_slide(prs, slide_data):
    """Build a single slide from JSON data."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    layout = slide_data.get('layout_type', 'content')
    title = slide_data.get('title', '')
    subtitle = slide_data.get('subtitle', '')
    content = slide_data.get('content', '')
    notes = slide_data.get('notes', '')
    style = slide_data.get('style', {})

    # Determine background
    is_dark = layout in ('title_slide', 'section')
    bg_info = style.get('background', {})
    if isinstance(bg_info, dict):
        bg_type = bg_info.get('type', '')
        if 'gradient' in bg_type or bg_info.get('color', '') in ('#1E2761', '#0A1240'):
            is_dark = True
        elif bg_info.get('color') == '#F0F4FF':
            set_slide_bg(slide, LIGHT_BG)

    if is_dark:
        set_slide_bg(slide, DARK_BLUE)

    # --- Title Slide ---
    if layout == 'title_slide':
        set_slide_bg(slide, DARK_BLUE)
        add_textbox(slide, 0.5, 2.0, 9, 1.5, title,
                   font_size=44, font_color=WHITE, bold=True,
                   alignment=PP_ALIGN.LEFT)
        if subtitle:
            add_textbox(slide, 0.5, 3.8, 9, 1.0, subtitle,
                       font_size=20, font_color=ICE_BLUE,
                       alignment=PP_ALIGN.LEFT)
        # Bottom accent line
        line = slide.shapes.add_shape(
            1, Inches(0.5), Inches(5.5), Inches(3), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_BLUE
        line.line.fill.background()

    # --- Section slide ---
    elif layout == 'section':
        set_slide_bg(slide, DARK_BLUE)
        add_textbox(slide, 0.5, 2.5, 9, 1.5, title,
                   font_size=40, font_color=WHITE, bold=True,
                   alignment=PP_ALIGN.CENTER)
        if subtitle:
            add_textbox(slide, 0.5, 4.2, 9, 1.0, subtitle,
                       font_size=18, font_color=ICE_BLUE,
                       alignment=PP_ALIGN.CENTER)

    # --- Code slide ---
    elif layout == 'code':
        add_title_bar(slide, title)
        items = flatten_content(content)
        # Code block background
        code_box = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = CODE_BG
        code_box.line.fill.background()

        code_text = '\n'.join(items)
        add_textbox(slide, 0.7, 1.6, 8.6, 5.0, code_text,
                   font_size=12, font_color=RGBColor(200, 220, 252),
                   font_name='Consolas')

    # --- Comparison / Table slide ---
    elif layout in ('comparison', 'table'):
        add_title_bar(slide, title)
        items = flatten_content(content)
        add_bullet_points(slide, 0.5, 1.4, 9, 5.5, items, font_size=13)

    # --- Steps slide ---
    elif layout == 'steps':
        add_title_bar(slide, title)
        items = flatten_content(content)
        y = 1.5
        for i, item in enumerate(items[:8], 1):
            # Step number circle
            circle = slide.shapes.add_shape(
                9, Inches(0.5), Inches(y), Inches(0.4), Inches(0.4))  # Oval
            circle.fill.solid()
            circle.fill.fore_color.rgb = DARK_BLUE
            circle.line.fill.background()
            # Number in circle
            tf = circle.text_frame
            tf.paragraphs[0].text = str(i)
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            add_textbox(slide, 1.1, y, 8.2, 0.5, str(item),
                       font_size=13, font_color=DARK_GRAY)
            y += 0.7

    # --- FAQ slide ---
    elif layout == 'faq':
        add_title_bar(slide, title)
        items = flatten_content(content)
        y = 1.5
        for item in items[:6]:
            # Q block with ice blue bg
            q_box = slide.shapes.add_shape(
                1, Inches(0.5), Inches(y), Inches(9), Inches(0.6))
            q_box.fill.solid()
            q_box.fill.fore_color.rgb = ICE_BLUE
            q_box.line.fill.background()
            add_textbox(slide, 0.7, y + 0.05, 8.6, 0.5, str(item),
                       font_size=12, font_color=DARK_BLUE)
            y += 0.8

    # --- Chart description slide ---
    elif layout == 'chart_description':
        add_title_bar(slide, title)
        items = flatten_content(content)
        # Placeholder for chart
        chart_area = slide.shapes.add_shape(
            1, Inches(0.5), Inches(1.4), Inches(5), Inches(5.0))
        chart_area.fill.solid()
        chart_area.fill.fore_color.rgb = LIGHT_BG
        chart_area.line.color.rgb = ICE_BLUE
        add_textbox(slide, 1.0, 3.0, 4, 1, "[圖表區域]",
                   font_size=16, font_color=DARK_GRAY,
                   alignment=PP_ALIGN.CENTER)
        # Data points on right
        if items:
            add_bullet_points(slide, 5.8, 1.4, 3.7, 5.0, items, font_size=12)

    # --- Architecture slide ---
    elif layout == 'architecture':
        add_title_bar(slide, title)
        items = flatten_content(content)
        add_bullet_points(slide, 0.5, 1.4, 9, 5.5, items, font_size=13)

    # --- Default content slide ---
    else:
        if is_dark:
            add_title_bar(slide, title, is_dark_bg=True)
            text_color = WHITE
            y_start = 1.6
        else:
            add_title_bar(slide, title)
            text_color = DARK_GRAY
            y_start = 1.4

        if subtitle and not is_dark:
            add_textbox(slide, 0.5, 1.2, 9, 0.5, subtitle,
                       font_size=16, font_color=ACCENT_BLUE)
            y_start = 1.8

        items = flatten_content(content)
        if items:
            add_bullet_points(slide, 0.5, y_start, 9, 5.5, items,
                            font_size=14, font_color=text_color)

    # Add notes
    add_notes(slide, notes)


def main():
    if len(sys.argv) != 3:
        print("Usage: python gen-pptx-from-json.py <input.json> <output.pptx>")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]

    # Ensure output directory exists
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    # Read JSON
    with open(input_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Handle both array and object-with-slides formats
    if isinstance(data, list):
        slides = data
    elif isinstance(data, dict):
        slides = data.get('slides', data.get('presentation', [data]))
        if isinstance(slides, dict):
            slides = [slides]
    else:
        print(f"Error: unexpected JSON format")
        sys.exit(1)

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Build slides
    for slide_data in slides:
        if not isinstance(slide_data, dict):
            continue
        try:
            build_slide(prs, slide_data)
        except Exception as e:
            print(f"Warning: Error building slide {slide_data.get('slide_number', '?')}: {e}")
            # Create a fallback slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_textbox(slide, 0.5, 3, 9, 1,
                       f"[Slide {slide_data.get('slide_number', '?')} - build error]",
                       font_size=20, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Save
    prs.save(output_path)
    print(f"Generated: {output_path} ({len(prs.slides)} slides)")


if __name__ == '__main__':
    main()
