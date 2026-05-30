#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RGBColor(30, 39, 97)      # #1E2761
ICE_BLUE = RGBColor(202, 220, 252)    # #CADCFC
WHITE = RGBColor(255, 255, 255)       # #FFFFFF
DARK_TEXT = RGBColor(51, 51, 51)      # #333333
LIGHT_GRAY = RGBColor(245, 245, 245)  # #F5F5F5
ACCENT = RGBColor(255, 107, 107)      # #FF6B6B

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with dark blue background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Add title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        sub_frame = sub_box.text_frame
        sub_p = sub_frame.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(24)
        sub_p.font.color.rgb = ICE_BLUE
        sub_p.font.name = 'Calibri'
        sub_p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_content_slide(prs, title):
    """Add a content slide with header bar"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Header bar
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_BLUE
    header_shape.line.color.rgb = DARK_BLUE
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.LEFT
    
    return slide

# Slide 1: Opening / Title Slide
add_title_slide(prs, "AI 虛擬組織課程", "打造智慧化的組織協作模式")

# Slide 2: 課程目標
slide = add_content_slide(prs, "課程目標")
objectives = [
    "了解 Agent 的概念與應用場景",
    "掌握虛擬組織的設計原則與實施方法",
    "學習如何衡量導入效益與投資回報",
    "建立組織內的 AI 智慧協作文化"
]
y_pos = 1.5
for obj in objectives:
    text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.5))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "• " + obj
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    y_pos += 0.6

# Slide 3: 課程架構
slide = add_content_slide(prs, "課程架構")
units = [
    ("第一單元", "Agent 基礎概念", "45分鐘"),
    ("第二單元", "虛擬組織的解法", "60分鐘"),
    ("第三單元", "實施與管理", "50分鐘"),
    ("第四單元", "效益試算與導入計畫", "35分鐘")
]

y_pos = 1.5
for idx, (phase, title, duration) in enumerate(units):
    # Unit box
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(y_pos), Inches(9), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ICE_BLUE if idx % 2 == 0 else LIGHT_GRAY
    shape.line.color.rgb = DARK_BLUE
    shape.line.width = Pt(2)
    
    # Phase
    phase_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.1), Inches(2.5), Inches(0.4))
    phase_frame = phase_box.text_frame
    p = phase_frame.paragraphs[0]
    p.text = phase
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = 'Calibri'
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.55), Inches(2.5), Inches(0.4))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    
    # Duration
    duration_box = slide.shapes.add_textbox(Inches(7.8), Inches(y_pos + 0.4), Inches(1.7), Inches(0.4))
    duration_frame = duration_box.text_frame
    p = duration_frame.paragraphs[0]
    p.text = duration
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.RIGHT
    
    y_pos += 1.35

# Slide 4: 第一單元 - Agent基礎概念
add_title_slide(prs, "第一單元", "Agent 基礎概念 (45分鐘)")

# Slide 5: 什麼是 Agent?
slide = add_content_slide(prs, "什麼是 Agent?")
intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
intro_frame = intro_box.text_frame
intro_frame.word_wrap = True
p = intro_frame.paragraphs[0]
p.text = "Agent 是一個能夠自主感知環境、進行決策和採取行動的智慧程式實體。"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = 'Calibri'

features = [
    "自主性 (Autonomy)：能獨立完成任務，無需人工干預",
    "感知能力 (Perception)：能理解並回應環境變化",
    "決策能力 (Decision-making)：基於規則或機器學習進行判斷",
    "行動能力 (Action)：執行具體任務並改變環境狀態"
]

y_pos = 2.8
for feature in features:
    text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.5))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "• " + feature
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    y_pos += 0.65

# Slide 6: Agent 的應用場景
slide = add_content_slide(prs, "Agent 的應用場景")
scenarios = [
    ("📧", "客服自動化", "24/7 自動應答客戶咨詢"),
    ("📊", "資料分析", "自動收集、處理和報表生成"),
    ("🔄", "流程自動化", "自動化業務流程和審批"),
    ("🎯", "決策支持", "實時數據分析和建議")
]

col_pos = 0.5
for icon, title, desc in scenarios:
    card_width = 2.1
    
    # Card background
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(col_pos), Inches(1.8), Inches(card_width), Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.color.rgb = ICE_BLUE
    shape.line.width = Pt(2)
    
    # Icon
    icon_box = slide.shapes.add_textbox(Inches(col_pos), Inches(2), Inches(card_width), Inches(0.6))
    icon_frame = icon_box.text_frame
    p = icon_frame.paragraphs[0]
    p.text = icon
    p.font.size = Pt(36)
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(col_pos + 0.1), Inches(2.8), Inches(card_width - 0.2), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(col_pos + 0.1), Inches(3.5), Inches(card_width - 0.2), Inches(1.5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    p = desc_frame.paragraphs[0]
    p.text = desc
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    col_pos += 2.3

# Slide 7: 第二單元 - 虛擬組織的解法
add_title_slide(prs, "第二單元", "虛擬組織的解法 (60分鐘)")

# Slide 8: 傳統組織的痛點
slide = add_content_slide(prs, "傳統組織的痛點")
pain_points = [
    "溝通低效：多層次審批導致響應遲緩",
    "協作困難：跨部門合作效率低下",
    "知識孤島：信息分散、難以共享",
    "人力成本高：大量重複性工作浪費人力資源"
]

y_pos = 1.8
for point in pain_points:
    # Red dot
    dot = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(y_pos), Inches(0.4), Inches(0.4)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.color.rgb = ACCENT
    
    # Text
    text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos), Inches(8.3), Inches(0.5))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = point
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    
    y_pos += 0.9

# Slide 9: Agent 組織的解法
slide = add_content_slide(prs, "Agent 組織的解法")
solutions = [
    ("任務複雜", "自動分派給對應專家", "⬇ 時間 -50%"),
    ("知識分散", "統一知識庫 + AI檢索", "⬆ 準確 +85%"),
    ("人工重複", "自動化常規操作", "⬇ 成本 -60%"),
    ("協作低效", "實時協作 + Agent聯動", "⬆ 效率 +70%")
]

y_pos = 1.7
for problem, solution, benefit in solutions:
    # Problem box
    problem_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(y_pos), Inches(2.3), Inches(0.6)
    )
    problem_shape.fill.solid()
    problem_shape.fill.fore_color.rgb = ACCENT
    problem_shape.line.color.rgb = DARK_BLUE
    
    problem_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.08), Inches(2.1), Inches(0.44))
    problem_frame = problem_box.text_frame
    p = problem_frame.paragraphs[0]
    p.text = problem
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    # Arrow
    arrow_box = slide.shapes.add_textbox(Inches(2.95), Inches(y_pos + 0.05), Inches(0.6), Inches(0.5))
    arrow_frame = arrow_box.text_frame
    p = arrow_frame.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Solution box
    solution_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(3.65), Inches(y_pos), Inches(3.2), Inches(0.6)
    )
    solution_shape.fill.solid()
    solution_shape.fill.fore_color.rgb = ICE_BLUE
    solution_shape.line.color.rgb = DARK_BLUE
    
    solution_box = slide.shapes.add_textbox(Inches(3.75), Inches(y_pos + 0.08), Inches(3), Inches(0.44))
    solution_frame = solution_box.text_frame
    p = solution_frame.paragraphs[0]
    p.text = solution
    p.font.size = Pt(12)
    p.font.color.rgb = DARK_BLUE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    # Benefit box
    benefit_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(6.95), Inches(y_pos), Inches(2.55), Inches(0.6)
    )
    benefit_shape.fill.solid()
    benefit_shape.fill.fore_color.rgb = LIGHT_GRAY
    benefit_shape.line.color.rgb = DARK_BLUE
    
    benefit_box = slide.shapes.add_textbox(Inches(7.05), Inches(y_pos + 0.08), Inches(2.35), Inches(0.44))
    benefit_frame = benefit_box.text_frame
    p = benefit_frame.paragraphs[0]
    p.text = benefit
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    y_pos += 0.8

# Slide 10: 虛擬組織架構
slide = add_content_slide(prs, "虛擬組織架構")
arch_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
arch_frame = arch_title.text_frame
p = arch_frame.paragraphs[0]
p.text = "多層次協作 Agent 系統"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = 'Calibri'

architecture = [
    ("協調層", "Manager Agent, Task Dispatcher"),
    ("執行層", "Specialist Agents (Sales, Support, Analytics, HR)"),
    ("支持層", "Knowledge Base, Data Bridge, Monitoring")
]

y_pos = 2.3
for idx, (layer, agents) in enumerate(architecture):
    bg_color = DARK_BLUE if idx == 0 else (ICE_BLUE if idx == 1 else LIGHT_GRAY)
    text_color = WHITE if idx == 0 else DARK_BLUE
    
    # Layer box
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(y_pos), Inches(9), Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = DARK_BLUE
    
    # Layer name
    layer_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.12), Inches(2), Inches(0.7))
    layer_frame = layer_box.text_frame
    p = layer_frame.paragraphs[0]
    p.text = layer
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    
    # Agents
    agents_box = slide.shapes.add_textbox(Inches(2.8), Inches(y_pos + 0.15), Inches(6.7), Inches(0.7))
    agents_frame = agents_box.text_frame
    p = agents_frame.paragraphs[0]
    p.text = agents
    p.font.size = Pt(12)
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    
    y_pos += 1.2

# Slide 11: 第三單元 - 實施與管理
add_title_slide(prs, "第三單元", "實施與管理 (50分鐘)")

# Slide 12: 實施路線圖
slide = add_content_slide(prs, "實施路線圖")
roadmap = [
    ("第1階段", "準備與評估", "2週"),
    ("第2階段", "試點與驗證", "4週"),
    ("第3階段", "正式部署", "2週"),
    ("第4階段", "優化與維護", "持續")
]

for idx, (phase, title, duration) in enumerate(roadmap):
    x_pos = 0.5 + (idx * 2.3)
    bg_color = DARK_BLUE if idx % 2 == 0 else ICE_BLUE
    text_color = WHITE if idx % 2 == 0 else DARK_BLUE
    
    # Phase box
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(x_pos), Inches(1.8), Inches(2), Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = DARK_BLUE
    
    # Phase label
    phase_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(1.9), Inches(1.8), Inches(0.3))
    phase_frame = phase_box.text_frame
    p = phase_frame.paragraphs[0]
    p.text = phase
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(2.25), Inches(1.8), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    # Duration
    duration_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(2.65), Inches(1.8), Inches(0.25))
    duration_frame = duration_box.text_frame
    p = duration_frame.paragraphs[0]
    p.text = duration
    p.font.size = Pt(10)
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

# Slide 13: Management Essentials
slide = add_content_slide(prs, "Management Essentials")
mgmt_points = [
    "1. 建立清晰的 Agent 角色定義與權限邊界",
    "2. 制定性能指標 (KPI) 監控 Agent 效能",
    "3. 實施數據安全與隱私保護措施",
    "4. 建立人類監督與干預機制",
    "5. 定期進行系統審計與優化"
]

y_pos = 1.8
for point in mgmt_points:
    text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.5))
    text_frame = text_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = point
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    y_pos += 0.7

# Slide 14: 第四單元 - 效益試算與導入計畫
add_title_slide(prs, "第四單元", "效益試算與導入計畫 (35分鐘)")

# Slide 15: 成本對比
slide = add_content_slide(prs, "成本對比：傳統 vs Agent 組織")

# Table header
header_shape = slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0.5), Inches(1.6), Inches(9), Inches(0.5)
)
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = DARK_BLUE
header_shape.line.color.rgb = DARK_BLUE

headers = ["項目", "傳統組織年成本", "Agent組織年成本"]
for idx, header in enumerate(headers):
    x_pos = 0.5 + (idx * 3)
    header_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.65), Inches(2.9), Inches(0.4))
    header_frame = header_box.text_frame
    p = header_frame.paragraphs[0]
    p.text = header
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

# Data rows
cost_data = [
    ("人力成本 (5人團隊)", "$250,000", "$85,000"),
    ("系統維護", "$15,000", "$8,000"),
    ("培訓與發展", "$12,000", "$20,000"),
    ("基礎設施", "$20,000", "$45,000")
]

y_pos = 2.2
for row_idx, (item, traditional, agent) in enumerate(cost_data):
    row_color = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
    
    # Row background
    row_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(y_pos), Inches(9), Inches(0.4)
    )
    row_shape.fill.solid()
    row_shape.fill.fore_color.rgb = row_color
    row_shape.line.color.rgb = ICE_BLUE
    
    # Item name
    item_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.02), Inches(2.8), Inches(0.36))
    item_frame = item_box.text_frame
    p = item_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    
    # Traditional cost
    trad_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos + 0.02), Inches(2.8), Inches(0.36))
    trad_frame = trad_box.text_frame
    p = trad_frame.paragraphs[0]
    p.text = traditional
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    # Agent cost
    agent_box = slide.shapes.add_textbox(Inches(6.5), Inches(y_pos + 0.02), Inches(2.8), Inches(0.36))
    agent_frame = agent_box.text_frame
    p = agent_frame.paragraphs[0]
    p.text = agent
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    y_pos += 0.5

# Total row
total_shape = slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5)
)
total_shape.fill.solid()
total_shape.fill.fore_color.rgb = ICE_BLUE
total_shape.line.color.rgb = DARK_BLUE

total_label_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.05), Inches(2.8), Inches(0.4))
total_label_frame = total_label_box.text_frame
p = total_label_frame.paragraphs[0]
p.text = "總計年成本"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = 'Calibri'

total_trad_box = slide.shapes.add_textbox(Inches(3.5), Inches(y_pos + 0.05), Inches(2.8), Inches(0.4))
total_trad_frame = total_trad_box.text_frame
p = total_trad_frame.paragraphs[0]
p.text = "$297,000"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DARK_TEXT
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

total_agent_box = slide.shapes.add_textbox(Inches(6.5), Inches(y_pos + 0.05), Inches(2.8), Inches(0.4))
total_agent_frame = total_agent_box.text_frame
p = total_agent_frame.paragraphs[0]
p.text = "$158,000"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

# Slide 16: ROI 分析
slide = add_content_slide(prs, "投資回報率 (ROI) 分析")
roi_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.4))
roi_title_frame = roi_title.text_frame
p = roi_title_frame.paragraphs[0]
p.text = "第一年投資與回報預測"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = 'Calibri'

roi_metrics = [
    ("初期投資", "$170,000", ACCENT),
    ("第一年成本節省", "$139,000", DARK_BLUE),
    ("生產效率提升", "$85,000", ICE_BLUE),
    ("總第一年效益", "$224,000", DARK_BLUE)
]

y_pos = 2.2
for metric, value, color in roi_metrics:
    metric_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(y_pos), Inches(9), Inches(0.5)
    )
    metric_shape.fill.solid()
    metric_shape.fill.fore_color.rgb = color
    metric_shape.line.color.rgb = DARK_BLUE
    
    text_color = WHITE if color in [ACCENT, DARK_BLUE] else DARK_BLUE
    
    metric_label_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.06), Inches(4.5), Inches(0.38))
    metric_label_frame = metric_label_box.text_frame
    p = metric_label_frame.paragraphs[0]
    p.text = metric
    p.font.size = Pt(13)
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    
    metric_value_box = slide.shapes.add_textbox(Inches(5.3), Inches(y_pos + 0.06), Inches(3.9), Inches(0.38))
    metric_value_frame = metric_value_box.text_frame
    p = metric_value_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.RIGHT
    
    y_pos += 0.65

# Slide 17: 實施案例
slide = add_content_slide(prs, "實施案例分析")

# Left box - Case study
case_shape = slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0.5), Inches(1.6), Inches(4.2), Inches(4.8)
)
case_shape.fill.solid()
case_shape.fill.fore_color.rgb = LIGHT_GRAY
case_shape.line.color.rgb = ICE_BLUE
case_shape.line.width = Pt(2)

case_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(3.8), Inches(0.4))
case_title_frame = case_title_box.text_frame
p = case_title_frame.paragraphs[0]
p.text = "案例：電商客服部門"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.font.name = 'Calibri'

case_points = [
    "背景：日均2000+客服詢問",
    "導入：3個 Agent 協同",
    "結果：",
    "  • 回應時間 ⬇ 85%",
    "  • 客戶滿意度 ⬆ 92%",
    "  • 成本節省 $180K/年"
]

y_pos = 2.4
for point in case_points:
    is_bold = not point.startswith("  ")
    font_size = 12 if is_bold and point not in ["結果："] else 11
    
    case_point_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(3.6), Inches(0.4))
    case_point_frame = case_point_box.text_frame
    p = case_point_frame.paragraphs[0]
    p.text = point
    p.font.size = Pt(font_size)
    p.font.bold = is_bold
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    y_pos += 0.5

# Right box - Metrics
metrics_shape = slide.shapes.add_shape(
    1,  # Rectangle
    Inches(5), Inches(1.6), Inches(4.2), Inches(4.8)
)
metrics_shape.fill.solid()
metrics_shape.fill.fore_color.rgb = DARK_BLUE
metrics_shape.line.color.rgb = ICE_BLUE
metrics_shape.line.width = Pt(2)

metrics_data = [
    ("投資回報期", "4 個月"),
    ("年度 ROI", "212%"),
    ("生產力提升", "3.5倍")
]

y_pos = 2.2
for label, value in metrics_data:
    metric_label_box = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos), Inches(3.8), Inches(0.4))
    metric_label_frame = metric_label_box.text_frame
    p = metric_label_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = ICE_BLUE
    p.font.name = 'Calibri'
    
    metric_value_box = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos + 0.5), Inches(3.8), Inches(0.6))
    metric_value_frame = metric_value_box.text_frame
    p = metric_value_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    
    y_pos += 1.4

# Slide 18: 導入建議
slide = add_content_slide(prs, "導入建議")
recommendations = [
    "1. 從小規模試點開始（選擇1-2個部門）",
    "2. 建立清晰的 Agent 角色與責任",
    "3. 投資員工培訓與文化建設",
    "4. 制定明確的績效指標與監測機制",
    "5. 預留應急資金應對不可預見的成本"
]

y_pos = 1.8
for rec in recommendations:
    rec_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8.5), Inches(0.5))
    rec_frame = rec_box.text_frame
    p = rec_frame.paragraphs[0]
    p.text = rec
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_TEXT
    p.font.name = 'Calibri'
    y_pos += 0.75

# Slide 19: 課程總結
closing_slide = add_title_slide(prs, "課程總結", "開啟智慧組織，創造無限可能")
closing_msg_box = closing_slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1))
closing_msg_frame = closing_msg_box.text_frame
closing_msg_frame.word_wrap = True
p = closing_msg_frame.paragraphs[0]
p.text = "AI 虛擬組織是未來企業協作的方向，透過合理的規劃與實施，可以顯著提升效率與競爭力。"
p.font.size = Pt(16)
p.font.color.rgb = WHITE
p.font.name = 'Calibri'
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "AI虛擬組織課程_改善版.pptx"
prs.save(output_path)
print(f"✅ PowerPoint 檔案已成功生成: {output_path}")
